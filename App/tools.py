# tools.py
import io
import json
import base64
import requests

from PIL import Image

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


# ==========================================
# Configuration Information
# ==========================================
import os
GEMINI_API_URL = os.getenv("GEMINI_API_URL", "YOUR_GEMINI_API_URL")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "YOUR_GEMINI_API_KEY")

# ==========================================
# Helper Functions: XML Operations (Not Exposed)
# ==========================================

def _apply_arrow_xml(connector, end_arrow=True, size='med'):
    """
    [Internal] Add arrow to line.
    Args:
        size: 'sm' (small), 'med' (medium), 'lg' (large)
    """
    ln = connector.line._get_or_add_ln()
    
    # 构建 XML: w=宽度, len=长度
    # 这里我们让长度和宽度保持一致
    arrow_xml = f'<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="{size}" len="{size}"/>'
    
    if not end_arrow: 
        arrow_xml = arrow_xml.replace('tailEnd', 'headEnd')
    
    ln.append(parse_xml(arrow_xml))

def _apply_gradient_xml(connector, color_start, color_end):
    """
    [关键修复版] 给连接线应用渐变色
    强制将 <a:gradFill> 插入到 XML 的第一个位置，防止被 PPT 视为损坏文件。
    """
    ln = connector.line._get_or_add_ln()
    
    # 1. 清理现有的填充属性 (solidFill, noFill 等)
    for child in list(ln):
        if 'Fill' in child.tag:
            ln.remove(child)
            
    # 2. 构造渐变 XML (红 -> 蓝, 左 -> 右)
    # 可以在这里扩展 ang (角度) 参数
    gradient_xml = f"""
    <a:gradFill {nsdecls('a')} flip="0" rotWithShape="1">
      <a:gsLst>
        <a:gs pos="0">
          <a:srgbClr val="{color_start}"/>
        </a:gs>
        <a:gs pos="100000">
          <a:srgbClr val="{color_end}"/>
        </a:gs>
      </a:gsLst>
      <a:lin ang="0" scaled="1"/>
    </a:gradFill>
    """
    
    # 3. 关键：插入到头部
    ln.insert(0, parse_xml(gradient_xml))

def _call_gemini_strict(prompt, aspect_ratio="1:1"):
    """Call Gemini API to generate image"""
    refined_prompt = (
        f"{prompt}. "
        "Vector icon style, flat design, minimal, "
        "solid white background, isolated, no shadows, "
        "clean edges, single object."
    )
    payload = json.dumps({
       "contents": [{"role": "user", "parts": [{"text": refined_prompt}]}],
       "generationConfig": {
          "responseModalities": ["IMAGE"],
          "imageConfig": {"aspectRatio": aspect_ratio}
       }
    })
    headers = {'Authorization': f'Bearer {GEMINI_API_KEY}', 'Content-Type': 'application/json'}

    try:
        response = requests.post(GEMINI_API_URL, headers=headers, data=payload)
        response.raise_for_status()
        response_json = response.json()
        if "candidates" in response_json and response_json["candidates"]:
            parts = response_json["candidates"][0]["content"]["parts"]
            for part in parts:
                if "inlineData" in part:
                    return io.BytesIO(base64.b64decode(part["inlineData"]["data"]))
        return None
    except Exception as e:
        print(f"❌ API Error: {e}")
        return None

def _remove_white_background(image_stream, tolerance=50, crop_tight=True):
    """Remove white background and crop"""
    try:
        img = Image.open(image_stream).convert("RGBA")
        datas = img.getdata()
        new_data = []
        for item in datas:
            if item[0] > (255 - tolerance) and item[1] > (255 - tolerance) and item[2] > (255 - tolerance):
                new_data.append((255, 255, 255, 0))
            else:
                new_data.append(item)
        img.putdata(new_data)
        
        if crop_tight:
            bbox = img.getbbox()
            if bbox: img = img.crop(bbox)
            
        output_stream = io.BytesIO()
        img.save(output_stream, format="PNG")
        output_stream.seek(0)
        return output_stream
    except Exception:
        image_stream.seek(0)
        return image_stream


def add_custom_route_arrow(slide, points, color='333333', width=2.0, end_arrow=True):
    """
    [Tool Function] Draw continuous polyline with custom path (Freeform).
    Fixed FreeformBuilder API call error.
    """
    if len(points) < 2:
        return None

    # 1. Create Builder (set starting point)
    # points[0] is (x, y) raw inch float
    start_x, start_y = points[0]
    builder = slide.shapes.build_freeform(Inches(start_x), Inches(start_y))
    
    # 2. Prepare subsequent node list
    # add_line_segments receives a list of coordinate tuples: [(emu_x, emu_y), ...]
    # We need to convert all remaining points to Inches (EMU) units
    vertex_pairs = [(Inches(x), Inches(y)) for x, y in points[1:]]
    
    # 3. Batch add line segments
    # close=False means do not close path (do not connect back to start), crucial for drawing arrows
    builder.add_line_segments(vertex_pairs, close=False)
        
    # 4. Convert to shape
    freeform_shape = builder.convert_to_shape()
    
    # 5. Style settings
    # Must clear fill, otherwise PPT will try to fill the area enclosed by this polyline
    freeform_shape.fill.background() 
    
    freeform_shape.line.width = Pt(width)
    freeform_shape.line.color.rgb = RGBColor(int(color[:2], 16), int(color[2:4], 16), int(color[4:], 16))
    
    # 6. Set arrow
    if end_arrow:
        _apply_arrow_xml(freeform_shape, end_arrow=True, size='lg')
        
    return freeform_shape


# ==========================================
# Tool 1: Smart Connection (Supports Gradient)
# ==========================================
def add_connector(slide, source_shape, dest_shape, type='curve', color='1F4E79', width=3.0, 
                  gradient_start=None, gradient_end=None,
                  arrow_size=None, 
                  conn_src=None, conn_dest=None,
                  ):
    
    # 1. Calculate center point (solve zero trap)
    s_x = source_shape.left + (source_shape.width // 2)
    s_y = source_shape.top + (source_shape.height // 2)
    d_x = dest_shape.left + (dest_shape.width // 2)
    d_y = dest_shape.top + (dest_shape.height // 2)

    # 2. Route selection
    if conn_src is not None and conn_dest is not None:
        src_idx = int(conn_src)
        dst_idx = int(conn_dest)
    else:
        # Auto routing (Auto Logic)
        # (3, 1) means Right->Left, so 3=Right, 1=Left
        dx, dy = d_x - s_x, d_y - s_y
        if abs(dx) > abs(dy):
            src_idx, dst_idx = (3, 1) if dx > 0 else (1, 3)
        else:
            src_idx, dst_idx = (2, 0) if dy > 0 else (0, 2)

    # 3. Create & Lock
    type_map = {'curve': MSO_CONNECTOR.CURVE, 'elbow': MSO_CONNECTOR.ELBOW, 'straight': MSO_CONNECTOR.STRAIGHT}
    conn_type = type_map.get(type, MSO_CONNECTOR.CURVE)
    
    connector = slide.shapes.add_connector(conn_type, s_x, s_y, d_x, d_y)
    connector.begin_connect(source_shape, src_idx)
    connector.end_connect(dest_shape, dst_idx)

    # 4. Style
    connector.line.width = Pt(width)
    if gradient_start and gradient_end:
        _apply_gradient_xml(connector, gradient_start.lstrip('#'), gradient_end.lstrip('#'))
    else:
        connector.line.color.rgb = _parse_color(color)

    # 5. Smart Arrow Size
    if arrow_size is None:
        # 如果线很细 (<2.0)，自动用小箭头
        if width < 2.0: final_size = 'sm'
        # 如果线中等 (2.0 - 4.5)，用中箭头
        elif width <= 4.5: final_size = 'med'
        # 如果线很粗 (>4.5)，用大箭头
        else: final_size = 'lg'
    else:
        final_size = arrow_size

    # 应用动态计算出的尺寸 (原代码这里是硬编码的 'lg')
    _apply_arrow_xml(connector, end_arrow=True, size=final_size)

    return connector

# ==========================================
# Tool 2: 自由坐标箭头 (无需 Shape)
# ==========================================

def add_free_arrow(slide, start_x, start_y, end_x, end_y, type='straight', color='1F4E79', width=3.0, gradient_start=None, gradient_end=None):
    """
    [Tool Function] 根据指定的坐标绘制箭头（不依赖 Shape 对象）。
    
    Args:
        slide: PPT slide 对象
        start_x, start_y: 起点坐标 (原始浮点数, 单位英寸)
        end_x, end_y: 终点坐标 (原始浮点数, 单位英寸)
        type: 'straight' (默认), 'curve', 'elbow'
        ... (其他样式参数同 add_connector)
    """
    
    # 1. 类型映射
    type_map = {
        'curve': MSO_CONNECTOR.CURVE,
        'elbow': MSO_CONNECTOR.ELBOW,
        'straight': MSO_CONNECTOR.STRAIGHT
    }
    conn_type_enum = type_map.get(type, MSO_CONNECTOR.STRAIGHT) # 自由绘制时默认由 straight 更直观

    # 2. 创建连接线
    # add_connector(type, begin_x, begin_y, end_x, end_y)
    # 注意：这里需要在内部进行 Inches 转换
    connector = slide.shapes.add_connector(
        conn_type_enum, 
        Inches(start_x), Inches(start_y), 
        Inches(end_x), Inches(end_y)
    )

    # 3. 设置样式 (线宽)
    connector.line.width = Pt(width)

    # 4. 设置颜色 (纯色 or 渐变)
    if gradient_start and gradient_end:
        _apply_gradient_xml(connector, gradient_start.lstrip('#'), gradient_end.lstrip('#'))
    else:
        color = color.lstrip('#')
        connector.line.color.rgb = RGBColor(int(color[:2], 16), int(color[2:4], 16), int(color[4:], 16))

    # 5. 加箭头
    _apply_arrow_xml(connector, end_arrow=True, size='lg')

    return connector
    


# ==========================================
# Tool 3: 自定义路径折线
# ==========================================
def add_custom_route_arrow(slide, points, color='333333', width=2.0, end_arrow=True):
    """
    [Tool Function] 绘制经过一系列自定义坐标点的折线箭头 (Freeform/Polyline)。
    用于解决 add_connector 无法实现的复杂绕行路径 (如 C 型、U 型连接)。

    Args:
        slide (obj): PPT slide 对象
        points (list of tuples): 坐标点列表 [(x1, y1), (x2, y2), ...]。单位为原始浮点数(英寸)。
        color (str): 线条颜色 Hex (e.g., 'FF0000')。
        width (float): 线宽 (points)。
        end_arrow (bool): 是否在最后一个点添加箭头。
    """
    # 至少需要两个点才能画线
    if not points or len(points) < 2:
        print("❌ Error: add_custom_route_arrow requires at least 2 points.")
        return None

    try:
        # 1. 创建 Freeform Builder (设置起点)
        # points[0] 是 (x, y) 原始英寸浮点数
        start_x, start_y = points[0]
        builder = slide.shapes.build_freeform(Inches(start_x), Inches(start_y))
        
        # 2. 准备后续节点列表 (转换为 EMU 单位)
        # builder.add_line_segments 需要 [(emu_x, emu_y), ...] 格式
        vertex_pairs = [(Inches(x), Inches(y)) for x, y in points[1:]]
        
        # 3. 批量添加线段
        # close=False 至关重要，否则会自动封闭成多边形
        builder.add_line_segments(vertex_pairs, close=False)
            
        # 4. 转换为形状
        freeform_shape = builder.convert_to_shape()
        
        # 5. 样式设置
        # 必须清除填充，确保它看起来是一条线而不是形状
        freeform_shape.fill.background() 
        
        freeform_shape.line.width = Pt(width)
        freeform_shape.line.color.rgb = RGBColor(int(color[:2], 16), int(color[2:4], 16), int(color[4:], 16))
        
        # 6. 设置箭头 (复用现有的 xml 函数)
        if end_arrow:
            _apply_arrow_xml(freeform_shape, end_arrow=True, size='lg')
            
        return freeform_shape

    except Exception as e:
        print(f"❌ Error in add_custom_route_arrow: {e}")
        return None
        


# ==============================================================================
# SECTION: UNIVERSAL DRAWING TOOLS (Block, Label, Container)
# ==============================================================================


# --- Helper: Color Parser ---
def _parse_color(color_input):
    """
    [Internal] Compatible handling of Hex string ('FF0000') or RGBColor object.
    Returns: RGBColor object
    """
    if isinstance(color_input, str):
        # 移除可能存在的 # 号
        hex_color = color_input.lstrip('#')
        # 防止空字符串导致报错，默认黑色
        if not hex_color:
            return RGBColor(0, 0, 0)
        return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))
    elif isinstance(color_input, RGBColor):
        return color_input
    # 默认 fallback
    return RGBColor(0, 0, 0)

# --- Helper: Transparency (Alpha) Injector ---
def _set_shape_alpha(shape, alpha):
    """
    [Internal] Set shape fill transparency via underlying XML.
    Args:
        shape: PPTX Shape object
        alpha: float, 0.0 (fully transparent) ~ 1.0 (opaque/Solid)
    """
    if alpha >= 1.0:
        return

    # 1. 获取 fill 属性
    fill = shape.fill
    # 必须先有 fore_color (即必须先设置为 solid fill) 才能设置 alpha
    if not hasattr(fill, 'fore_color'):
        return

    # 2. 获取底层的 srgbClr 元素 (Solid Color)
    try:
        # python-pptx 的 fill.fore_color._xFill 对应 <a:solidFill>
        xFill = fill.fore_color._xFill 
        # 获取 <a:srgbClr>
        srgbClr = xFill.srgbClr
        
        if srgbClr is not None:
            # 先移除旧的 alpha 节点 (如果有，防止重复叠加)
            alpha_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
            existing_alphas = srgbClr.findall(alpha_tag)
            for a in existing_alphas:
                srgbClr.remove(a)
            
            # 转换为 XML 整数: 100% = 100000
            alpha_int = int(alpha * 100000)
            
            # 注入新的 alpha 节点
            alpha_elem = parse_xml(f'<a:alpha {nsdecls("a")} val="{alpha_int}"/>')
            srgbClr.append(alpha_elem)
                
    except Exception as e:
        print(f"⚠️ Warning: Failed to set transparency alpha={alpha}: {e}")

# --- Tool 4: Universal Shape Creator (add_block) ---
def add_block(slide, x, y, w, h, text=None, fill_color='FFFFFF', stroke_color='000000', 
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, font_size=12, font_color='000000', bold=False,
              alpha=1.0):
    """
    [Tool] Draw a generic entity block (node).
    Supports: auto-wrap, center alignment, border, safe fill, transparency.
    """
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    
    # 1. 填充 (Fill) & 透明度 (Alpha)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(fill_color)
        if alpha < 1.0:
            _set_shape_alpha(shape, alpha)
    else:
        shape.fill.background()

    # 2. 边框 (Line)
    if stroke_color:
        shape.line.color.rgb = _parse_color(stroke_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background() # 无边框

    # 3. 文本 (Text)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE # 垂直居中
        
        # Get or create paragraph
        if len(tf.paragraphs) == 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
            
        p.text = str(text)
        p.alignment = PP_ALIGN.CENTER # 水平居中
        p.font.size = Pt(font_size)
        p.font.color.rgb = _parse_color(font_color)
        p.font.bold = bold
    
    return shape

# --- Tool 5: Pure Text Label (add_label) ---
def add_label(slide, text, x, y, w=None, h=None, font_size=12, color='000000', bold=False, align='center'):
    """
    [Tool] Draw pure text label (no fill, no border).
    If w or h is None, enable 'SHAPE_TO_FIT_TEXT' auto-resize.
    """
    # Default size (if not specified and not auto-adjusting)
    width_in = Inches(w) if w else Inches(2.0)
    height_in = Inches(h) if h else Inches(0.5)
    
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), width_in, height_in)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    # If width/height not specified, enable auto-fit
    if w is None: 
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Set text content
    if len(tf.paragraphs) == 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = _parse_color(color)
    p.font.bold = bold

    # Alignment mapping
    align_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    p.alignment = align_map.get(str(align).lower(), PP_ALIGN.CENTER)
    
    return textbox

# --- Tool 6: Container/Group Frame (add_container) ---
def add_container(slide, x, y, w, h, title=None, fill_color='F5F5F5', stroke_color='CCCCCC', alpha=1.0):
    """
    [Tool] Draw a background container, usually used for grouping.
    Supports transparency (alpha). If title is provided, will draw title at top inside container.
    """
    # 1. Draw background rectangle
    container = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    
    # 2. 填充 & 透明度
    if fill_color:
        container.fill.solid()
        container.fill.fore_color.rgb = _parse_color(fill_color)
        if alpha < 1.0:
            _set_shape_alpha(container, alpha)
    else:
        container.fill.background()
        
    # 3. 边框
    if stroke_color:
        container.line.color.rgb = _parse_color(stroke_color)
        container.line.width = Pt(1.5)
    else:
        container.line.fill.background()
        
    # 4. Remove shadow (flat style)
    container.shadow.inherit = False

    # 5. Optional title (call add_label)
    if title:
        # Title located at container top, height fixed at 0.4 inches
        add_label(
            slide, title, 
            x=x, y=y + 0.1, w=w, h=0.4, 
            font_size=11, bold=True, color='333333', align='center'
        )
        
    return container