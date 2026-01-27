# ==============================================================================
# å®šä¹‰ python-pptx çš„â€œé¿å‘æŒ‡å—â€
# å°†è¿™äº›è§„åˆ™æ³¨å…¥åˆ° Prompt ä¸­ï¼Œèƒ½å¤Ÿå¤§å¹…å‡å°‘ 90% çš„å¸¸è§æŠ¥é”™
# ==============================================================================
PPTX_BEST_PRACTICES = """
*** CRITICAL PYTHON-PPTX RULES (MUST FOLLOW) ***
1. **Lines are CONNECTORS**: 
   - NEVER use `slide.shapes.add_shape(MSO_SHAPE.LINE, ...)` -> This causes AttributeError.
   - ALWAYS use `slide.shapes.add_connector(MSO_CONNECTOR.X, ...)`.
   - **Valid Types**: `MSO_CONNECTOR.STRAIGHT`, `MSO_CONNECTOR.ELBOW`, `MSO_CONNECTOR.CURVE`.
   - **INVALID**: Do NOT use `MSO_CONNECTOR.CURVED` (No 'D' at the end).
   - **INVALID SHAPES**: `MSO_SHAPE.DOC_TAG` does not exist (Use `MSO_SHAPE.FOLDED_CORNER` instead).
   
2. **Connector Properties**:
   - Connectors (Lines/Arrows) have `.line` but **NO `.fill`**. 
   - NEVER try to set `connector.fill.solid()`. Only set `connector.line.color.rgb`.

3. **Shape Fills (NO ONE-LINERS)**:
   - **NEVER** try to create and color a shape in one line: `add_shape(...).fill.fore_color.rgb = ...` (This crashes with TypeError).
   - **ALWAYS** split into steps:
     1. `shape = slide.shapes.add_shape(...)`
     2. `shape.fill.solid()`  <-- REQUIRED first!
     3. `shape.fill.fore_color.rgb = RGBColor(...)`

4. **Imports & Enums (STRICT)**:
   - Use these exact imports to avoid ImportErrors:
     ```python
     from pptx import Presentation
     from pptx.util import Inches, Pt, Cm
     from pptx.dml.color import RGBColor
     from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
     from pptx.enum.dml import MSO_LINE_DASH_STYLE
     from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
     ```
   - **FORBIDDEN IMPORTS (DO NOT USE)**:
     - `MSO_ARROWHEAD` (Does not exist. Use tools for arrows).
     - `MSO_TEXT_ORIENTATION` (Causes ImportError. Do not rotate text via enum).
   - **Line Dashes**: Use `line.dash_style = MSO_LINE_DASH_STYLE.DASH`.

5. **Text Frames**:
   - Always check `if shape.has_text_frame:` before accessing text properties.

6. **Text Styling & Helper Functions**:
   - If you define a helper function like `def add_text(...)`, **DO NOT** use hypothetical arguments like `style='italic'` or `css='bold'` when calling it.
   - **Italic**: Use `run.font.italic = True`.
   - **Bold**: Use `run.font.bold = True`.
   - **Underline**: Use `run.font.underline = True`.
   - Do NOT pass "style" strings. Set attributes explicitly on the font object.

7. **Color Assignment Safety (CRITICAL)**:
   - The `.rgb` property ONLY accepts **`RGBColor` objects**.
   - It **REJECTS** `None`, `tuples` (e.g. `(255,0,0)`), and `hex strings`.
   - **Correct Pattern**:
     ```python
     # 1. Check for None
     if color_variable is not None:
         # 2. Ensure it is an RGBColor object (not a tuple)
         shape.fill.fore_color.rgb = color_variable 
     ```
   - If you have a tuple `(r, g, b)`, wrap it: `RGBColor(r, g, b)`.

8. **Coordinate Precision & Type Safety (CRITICAL)**:
   - **Integers Only**: Coordinates (x,y,w,h) MUST be integers (EMUs).
   - Python division (`/`) produces floats (e.g. `Cm(10)/2` -> float). Floats cause crashes in `python-pptx` XML generation.
   - **FIX**: Always wrap calculations in `int()`: `int(Cm(4) * 1.5)`.
   - **No Reverse Access**: The result of `Inches(1) + Inches(2)` is a plain `int`.
   - **NEVER** try to access `.inches`, `.cm`, or `.pt` on a coordinate variable (e.g., `x.inches` -> AttributeError).

9. **Tools Usage**:
   - Check if you need to draw arrows. If yes, `from tools import add_connector`.
   - **Do not redefine** these functions. They are already available in `tools.py`.
"""

# PPTX_BEST_PRACTICES = ""




# ==============================================================================
# å®šä¹‰ Tools çš„ API è¯´æ˜Ž (ä½œä¸º Prompt çš„ä¸€éƒ¨åˆ†)
# ==============================================================================
# ==============================================================================
# TOOLS_SPECIFICATION: å®šä¹‰ Tools çš„ API è¯´æ˜Ž (ä½œä¸º Prompt çš„ä¸€éƒ¨åˆ†)
# ==============================================================================
TOOLS_SPECIFICATION = """
*** HIGH-PRIORITY TOOLS SPECIFICATION ***

You have access to a local library `tools.py` that provides high-level plotting capabilities.
**You must always prioritize using these tools over native `python-pptx` methods to ensure best visual quality.**

---

#### **âš ï¸ Global Rules & Best Practices**

1. **Imports**: **ALWAYS** use wildcard import to get all tools:
```python
   from tools import *
```

2. **Coordinate Units**:
* For **Tools** (e.g., `add_block`, `add_connector`): Use **raw floats** (e.g., `left=5.0`).
* For **Native PPTX** (`slide.shapes.add_shape`): Use **`Inches()`** (e.g., `left=Inches(5.0)`).

3. **Routing**: Do not calculate connection indices manually. The tools handle alignment automatically.
4. **Objects**: Always pass Shape/Picture objects to connector functions (`add_connector`), not their names.
5. **Strict Parameter Compliance**: The function signatures listed below are EXHAUSTIVE. DO NOT use any parameters that are not explicitly defined in the documentation (e.g., do not hallucinate linestyle, dashed, shadow, or end_arrow unless they appear in the signature).

---

### **ðŸ“¦ SECTION 1: UNIVERSAL DRAWING TOOLS (Nodes, Text, Groups)**

#### **Tool 1: `add_container` (Background Grouping)**

**Description**
Draws a background rectangle to visually group related elements.

* **Best Practice**: Call this **FIRST** (Layer 1) before drawing nodes inside it, to ensure it stays in the background.

**Function Signature**

```python
add_container(slide, x, y, w, h, title=None, fill_color='F5F5F5', stroke_color='CCCCCC', alpha=1.0)

```

**Parameters**

* **`title`** *(str, optional)*: Automatically adds a bold title at the top inside the container.
* **`alpha`** *(float)*: Transparency. `1.0` is opaque, `0.0` is invisible. Use `0.1`-`0.3` for subtle backgrounds.

**Example**

```python
# Draw a light grey background area for the "Encoder" section
group_box = add_container(slide, x=0.5, y=1.0, w=4.0, h=5.0, title="Encoder Layers")

```

#### **Tool 2: `add_block` (Generic Node/Shape)**

**Description**
The **primary tool** for drawing nodes, boxes, steps, or components. Handles text centering, safe color filling, borders, and transparency automatically.

**Function Signature**

```python
add_block(slide, x, y, w, h, text=None, fill_color='FFFFFF', stroke_color='000000', 
          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, font_size=12, font_color='000000', bold=False, alpha=1.0)

```

**Parameters**

* **`fill_color`** *(str)*: Hex code (e.g., `'FF0000'`).
* **`alpha`** *(float)*: Transparency. Use `alpha=0.5` for overlays/highlighting overlapping regions.
* **`shape_type`**: The specific geometry of the node. **STRICTLY** use one of the following constants:
    * **Rectangles**: `MSO_SHAPE.RECTANGLE`, `MSO_SHAPE.ROUNDED_RECTANGLE`, `MSO_SHAPE.SNIP_1_RECTANGLE` (Folded Corner)
    * **Circles/Ovals**: `MSO_SHAPE.OVAL` (Use equal w/h for Circle)
    * **Polygons**: `MSO_SHAPE.DIAMOND`, `MSO_SHAPE.ISOSCELES_TRIANGLE`, `MSO_SHAPE.HEXAGON`, `MSO_SHAPE.PENTAGON`
    * **Arrows**: `MSO_SHAPE.RIGHT_ARROW`, `MSO_SHAPE.LEFT_ARROW`, `MSO_SHAPE.UP_ARROW`, `MSO_SHAPE.DOWN_ARROW`
    * **Others**: `MSO_SHAPE.CHEVRON`, `MSO_SHAPE.CLOUD`, `MSO_SHAPE.HEART`, `MSO_SHAPE.TEAR`

**Example**

```python
# Draw a standard blue node
node = add_block(slide, x=1.0, y=2.0, w=3.0, h=1.0, text="Transformer Block", fill_color='DAE8FC')

# Draw a semi-transparent red overlay
overlay = add_block(slide, x=1.5, y=2.5, w=1.0, h=1.0, text="Focus", fill_color='FF0000', alpha=0.4)

```

#### **Tool 3: `add_label` (Pure Text)**

**Description**
Draws text without background or borders. Use for titles, captions, or annotations floating outside blocks.

**Function Signature**

```python
add_label(slide, text, x, y, w=None, h=None, font_size=12, color='000000', bold=False, align='center')

```

**Parameters**

* **`w`** *(float, optional)*: If `None`, the text box will **auto-resize** to fit the content (`SHAPE_TO_FIT_TEXT`).
* **`align`**: `'left'`, `'center'`, `'right'`.

**Example**

```python
# Add a main title
add_label(slide, "Figure 1: Architecture", x=0.5, y=0.2, font_size=18, bold=True)

```

---

### **ðŸ”— SECTION 2: CONNECTION TOOLS (Arrows & Lines)**


#### **Tool 4: `add_connector` (Smart Arrow)**

**Description**
Draws a connecting line with an automatic arrowhead between two existing shape objects. Supports intelligent path routing and **gradient color effects**.

**Function Signature**

```python
add_connector(slide, source_shape, dest_shape, type='curve', color='1F4E79', width=3.0, gradient_start=None, gradient_end=None, arrow_size=None)

```

**Parameters**

* **`slide`** *(obj)*: The active slide object.
* **`source_shape`**, **`dest_shape`** *(obj)*: The start and end Shape/Picture objects (returned by `add_block`).
* **`type`** *(str)*:
* `'curve'` (Default): S-curve. Best for horizontal flows.
* `'elbow'`: Orthogonal lines (90-degree turns). Best for vertical hierarchies.
* `'straight'`: Direct straight line.


* **`width`** *(float)*: Line thickness in points. Default `3.0`.
* **`color`** *(str)*: Hex code for solid color (e.g., `'FF0000'`). Ignored if gradients are used.
* **`gradient_start`** *(str, optional)*: Hex code for gradient start (e.g., `'FF0000'`).
* **`gradient_end`** *(str, optional)*: Hex code for gradient end (e.g., `'0000FF'`).
* **`arrow_size`** *(str, optional)*: Controls the size of the arrowhead.
    * Options: `'sm'` (Small), `'med'` (Medium), `'lg'` (Large).
    * **Default**: `None` (Auto-calculated based on line width: Thin lines get small arrows).

**âš ï¸ Limitations & Best Use**

* Best For: Simple, direct connections where the line can go straight from A to B (or via a simple elbow) without hitting obstacles.
* Routing Logic: This tool optimizes for the shortest/most direct path. It CANNOT draw complex routes, bypass obstacles, or pass through specific waypoints.
* Alternative: If you need a specific path (e.g., C-shape, U-turn, or avoiding overlap), DO NOT use this tool. Use add_custom_route_arrow instead.

**Usage Logic**

* **Solid Line**: Provide `color`.
* **Gradient Line**: Provide both `gradient_start` and `gradient_end`.

**Example**

```python
# 1. Standard Solid Curve
add_connector(slide, node_a, node_b, type='curve', color='1F4E79')

# 2. Gradient Flow (e.g., Input to Output)
add_connector(slide, node_input, node_output, type='curve', width=5.0, 
              gradient_start='FF7F50', gradient_end='8A2BE2')

```


#### **Tool 5: `add_free_arrow` (Coordinate-Based Arrow)**

**Description**
Draws an arrow between two arbitrary (x, y) coordinates.

* **Use Case**: Use this when you need to point to a specific location, draw axis lines, or create annotations **that are NOT attached to specific shape objects**.
* **Contrast**: Use `add_connector` if you want the line to stay attached when shapes move. Use `add_free_arrow` for absolute positioning.

**Function Signature**

```python
add_free_arrow(slide, start_x, start_y, end_x, end_y, type='straight', color='1F4E79', width=3.0, gradient_start=None, gradient_end=None)

```

**Parameters**

* **`slide`** *(obj)*: The active slide object.
* **`start_x`, `start_y**` *(float)*: Starting coordinates in **raw inches**.
* **`end_x`, `end_y**` *(float)*: Ending coordinates in **raw inches**.
* **`type`** *(str)*: `'straight'` (Default), `'curve'`, `'elbow'`.
* *(Style parameters `color`, `width`, `gradient_*` are same as `add_connector`)*

**Example**

```python
# 1. Annotation: Pointing from a text label (3.0, 2.0) to a detail (5.5, 4.0)
add_free_arrow(slide, 3.0, 2.0, 5.5, 4.0, type='straight', color='FF0000')

# 2. Divider Line: Drawing a horizontal gradient separator at the bottom
add_free_arrow(slide, 1.0, 7.0, 12.0, 7.0, width=4.0, gradient_start='000000', gradient_end='CCCCCC')

```

#### **Tool 6: `add_custom_route_arrow` (Custom Polyline Path)**

**Description**
Draws a multi-segment arrow passing through a specific sequence of `(x, y)` coordinates.

* **Purpose**: Use this tool when standard `add_connector` fails to route correctly (e.g., when you need to route around an obstacle, create a C-shaped loop, or draw a U-turn).
* **Mechanism**: It creates a single continuous "Freeform" shape (Polyline) defined exactly by the vertices you provide.

**Function Signature**

```python
add_custom_route_arrow(slide, points, color='333333', width=2.0, end_arrow=True)

```

**Parameters**

* **`points`** *(list of tuples)*: A strictly ordered list of `(x, y)` float coordinates defining the turning points.
* Format: `[(x1, y1), (x2, y2), (x3, y3), ...]`
* **Unit**: Raw Inches (float). Do NOT use `Inches()`.
* **`color`** *(str)*: Hex color code (e.g., `'FF0000'`).
* **`width`** *(float)*: Line thickness in points.
* **`end_arrow`** *(bool)*: Whether to add an arrow head at the last point. Default is `True`.

**ðŸ§  How to Plan the Path (The "Lane" Strategy)**

When connecting Shape A to Shape B with a complex path, do not guess random coordinates. Use the **"Lane Strategy"**:

1. **Identify Anchors**: Get the coordinates of the Start (`A`) and End (`B`).
2. **Define a Lane**: Choose an empty X or Y coordinate to act as a "highway" for the line (e.g., "0.5 inches to the left of A").
3. **Construct Segments**: Build the list of points to move from Start -> Lane -> End.

**Common Path Patterns:**

* **L-Shape (Simple Turn)**: `[Start, Corner, End]`
* **C-Shape (Bypass/Residual)**: `[Start, Move_Out, Move_Along_Lane, Move_In]`
* **U-Shape (Feedback)**: `[Start_Bottom, Move_Right, Move_Up, End_Top]`

**Usage Example**

**Scenario**: Drawing a "Residual Connection" (Skip Connection) that goes from the *Bottom Shape* to the *Top Shape*, bypassing a *Middle Shape* by routing around the **Left** side.

```python
# 1. Get Geometry
# Start at the LEFT side of the Bottom Node
x_start = node_bottom.left.inches
y_start = node_bottom.top.inches + (node_bottom.height.inches / 2)

# End at the LEFT side of the Top Node
x_end = node_top.left.inches
y_end = node_top.top.inches + (node_top.height.inches / 2)

# 2. Define the "Lane"
# Calculate a vertical track 0.6 inches to the left of the Bottom Node
x_lane = node_bottom.left.inches - 0.6

# 3. Construct the Points Sequence
# Path: Start -> Move Left to Lane -> Move Up along Lane -> Move Right to End
points_path = [
    (x_start, y_start),   # Point 1: Origin
    (x_lane, y_start),    # Point 2: Enter Lane
    (x_lane, y_end),      # Point 3: Travel Up Lane
    (x_end, y_end)        # Point 4: Exit to Destination
]

# 4. Draw
add_custom_route_arrow(slide, points=points_path, color='333333', width=2.5)

```
"""

# TOOLS_SPECIFICATION = ""